from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .config import FIGURES_DIR, METRICS_DIR, REPORTS_DIR, STUDENT_ID, STUDENT_NAME


INK = RGBColor(18, 38, 58)
TEAL = RGBColor(21, 122, 110)
GOLD = RGBColor(192, 139, 48)
SAND = RGBColor(246, 240, 232)
WHITE = RGBColor(255, 255, 255)


def _set_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = SAND


def _add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.0), Inches(0.9))
    title_frame = title_box.text_frame
    title_frame.clear()
    paragraph = title_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = INK
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(12.0), Inches(0.45))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.clear()
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(13)
        p.font.color.rgb = TEAL


def _add_bullets(slide, items: list[str], left: float = 0.8, top: float = 1.6, width: float = 5.6, height: float = 4.8) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    for idx, item in enumerate(items):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = INK
        paragraph.space_after = Pt(8)


def _add_picture(slide, filename: str, left: float, top: float, width: float) -> None:
    slide.shapes.add_picture(str(FIGURES_DIR / filename), Inches(left), Inches(top), width=Inches(width))


def build_presentation() -> None:
    metrics = pd.read_csv(METRICS_DIR / "model_comparison_metrics.csv")
    best = metrics.sort_values("f1", ascending=False).iloc[0]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.55), Inches(12.2), Inches(2.0))
    banner.fill.solid()
    banner.fill.fore_color.rgb = WHITE
    banner.line.color.rgb = GOLD
    _add_title(slide, "Xente Loan Default Prediction", f"Prepared by {STUDENT_NAME} ({STUDENT_ID})")
    _add_bullets(
        slide,
        [
            "Business, Management and Financial Data Analytics",
            "Leakage-aware machine learning workflow",
            "Jupyter notebooks, premium visuals, and Streamlit deployment",
        ],
        left=0.9,
        top=1.8,
        width=6.2,
        height=2.7,
    )
    _add_picture(slide, "model_comparison_metrics.png", left=7.4, top=1.55, width=5.1)

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "Business Problem and Objective")
    _add_bullets(
        slide,
        [
            "Xente extends digital credit through a mobile-first financial platform.",
            "Loan default erodes profitability, liquidity, and long-term sustainability.",
            "Objective: predict whether a borrower is likely to default on a loan.",
            "Goal: support earlier and more consistent credit decisions.",
        ],
    )
    _add_picture(slide, "target_distribution.png", left=7.1, top=1.45, width=5.6)

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "Dataset and Modelling Unit")
    _add_bullets(
        slide,
        [
            "TRAIN.csv used for EDA, cleaning, feature engineering, and validation.",
            "TEST.csv reserved for final unseen scoring.",
            "unlinked_masked_final.csv used to enrich prior customer behaviour.",
            "A practical row-level model was used while acknowledging repeated customer and loan events.",
        ],
    )
    _add_picture(slide, "monthly_trends.png", left=7.0, top=1.45, width=5.7)

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "Data Quality and Leakage Control")
    _add_bullets(
        slide,
        [
            "Only labelled observations were used for supervised learning.",
            "Missingness was profiled column-by-column instead of handled blindly.",
            "Outliers were treated with business-aware log scaling and IQR clipping.",
            "Repayment and payback fields were excluded to prevent leakage.",
        ],
        width=5.8,
    )
    _add_picture(slide, "missingness_top_columns.png", left=7.0, top=1.45, width=5.7)

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "EDA Highlights")
    _add_picture(slide, "distribution_Amount.png", left=0.6, top=1.4, width=6.0)
    _add_picture(slide, "categories_ProductCategory.png", left=6.8, top=1.4, width=5.8)

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "Feature Engineering and Selection")
    _add_bullets(
        slide,
        [
            "Transaction intensity: amount, value, ratios, and log-scaled variants.",
            "Temporal context: month, day, hour, weekday, and weekend indicator.",
            "Behavioural history: prior counts, prior spending, diversity, recency, and tenure.",
            "Raw IDs were excluded; behaviour summaries were retained.",
        ],
        width=5.8,
    )
    _add_picture(slide, "relationship_categorical_ProductCategory.png", left=7.0, top=1.45, width=5.7)

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "Model Comparison")
    _add_picture(slide, "model_comparison_metrics.png", left=0.6, top=1.4, width=6.1)
    _add_picture(slide, "roc_curve_comparison.png", left=6.9, top=1.4, width=5.8)

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "Best Model and Threshold")
    _add_bullets(
        slide,
        [
            f"Best model: {best['model']}",
            f"F1-score: {best['f1']:.3f}",
            f"ROC-AUC: {best['roc_auc']:.3f}",
            f"Recommended threshold: {best['threshold']:.2f}",
            "Threshold tuning improved the balance between precision and recall.",
        ],
        width=5.4,
    )
    _add_picture(slide, "threshold_tradeoff.png", left=6.5, top=1.5, width=6.0)

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "Most Influential Features")
    _add_picture(slide, "best_model_feature_effects.png", left=0.8, top=1.4, width=11.6)

    # Slide 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "Business Recommendations")
    _add_bullets(
        slide,
        [
            "Use the score as an early warning signal in digital credit approval.",
            "Flag medium-risk and high-risk cases for manual review.",
            "Combine the model with business rules around product usage and borrower history.",
            "Monitor model drift and recalibrate periodically.",
        ],
    )
    _add_picture(slide, "confusion_matrix_random_forest.png", left=7.1, top=1.45, width=5.5)

    # Slide 11
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_title(slide, "Deployment")
    _add_bullets(
        slide,
        [
            "Streamlit app includes data overview, EDA, cleaning logic, model performance, and live scoring.",
            "The app loads saved artifacts instead of retraining on every run.",
            "Notebook deliverables document the full analytical workflow.",
        ],
        width=5.7,
    )
    _add_picture(slide, "model_comparison_metrics.png", left=7.0, top=1.5, width=5.6)

    output_path = REPORTS_DIR / f"Xente_Loan_Default_Presentation_{STUDENT_ID}.pptx"
    prs.save(output_path)


if __name__ == "__main__":
    build_presentation()

